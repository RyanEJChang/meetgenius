import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pathlib import Path

# 前端 accept 屬性與後端驗證都應以此為準
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.md', '.txt', '.csv'}

# 純文字類格式：直接讀取即可
_PLAIN_TEXT_EXTENSIONS = {'.md', '.txt', '.csv'}


def read_document_text(file_path: str) -> str:
    """
    讀取文件的文字內容。

    支援 PDF、Word (.docx)、簡報 (.pptx) 與純文字類 (.md / .txt / .csv)。

    Args:
        file_path: 文件路徑。

    Returns:
        擷取出的文字內容。

    Raises:
        ValueError: 不支援的檔案格式。
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    text_content = ""

    if suffix == '.pdf':
        with fitz.open(path) as doc:
            for page in doc:
                text_content += page.get_text()
    elif suffix == '.docx':
        doc = Document(path)
        for para in doc.paragraphs:
            text_content += para.text + '\n'
        # 表格常放議程、負責人與時程，不能漏
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    text_content += ' | '.join(cells) + '\n'
    elif suffix == '.pptx':
        pres = Presentation(path)
        for index, slide in enumerate(pres.slides, 1):
            text_content += f'\n[投影片 {index}]\n'
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content += shape.text + '\n'
    elif suffix in _PLAIN_TEXT_EXTENSIONS:
        # 以 UTF-8 讀取，遇到少數無法解碼的位元組不中斷
        text_content = path.read_text(encoding='utf-8', errors='replace')
    else:
        raise ValueError(f"Unsupported file format: {suffix}")

    return text_content.strip()


def read_reference_document(file_path: str, max_chars: int = 6000) -> str:
    """
    讀取參考文件供摘要使用，失敗時回傳空字串而非拋出例外。

    參考文件只是輔助資訊，解析失敗不應該讓整場會議的處理中斷。

    Args:
        file_path: 文件路徑。
        max_chars: 最多取用的字元數，避免壓過逐字稿的比重。
    """
    if not file_path:
        return ""

    path = Path(file_path)
    if not path.exists():
        print(f"[WARNING] 參考文件不存在，略過：{file_path}")
        return ""

    try:
        text = read_document_text(str(path))
    except Exception as e:
        print(f"[WARNING] 參考文件解析失敗，略過：{path.name}（{e}）")
        return ""

    if not text:
        print(f"[WARNING] 參考文件沒有可讀取的文字，略過：{path.name}")
        return ""

    if len(text) > max_chars:
        print(f"[INFO] 參考文件過長（{len(text)} 字元），截斷至 {max_chars} 字元")
        text = text[:max_chars] + "\n…（後續內容已截斷）"

    return text
